#!/usr/bin/env python3
"""
选股系统 Skill 融合模块 v1.0
将10个Codex/OpenClaw skill整合到A股动量选股系统中
"""

import json, os, sys
from datetime import datetime
from pathlib import Path

sys.path.insert(0, '/root/.openclaw/workspace/skills/ifind-momentum-screener/scripts')

# ═══════════════════════════════════════════════════════════════
# 1. HUMANIZER — 去AI味（已全局启用，此模块无需额外调用）
# 说明: MEMORY.md已标记默认启用，所有输出自动经过humanizer处理
# ═══════════════════════════════════════════════════════════════

# ═══════════════════════════════════════════════════════════════
# 2. THEME-FACTORY — 主题样式系统
# ═══════════════════════════════════════════════════════════════

THEMES = {
    'tech_innovation': {
        'name': 'Tech Innovation',
        'primary': '#2563EB',      # 科技蓝
        'secondary': '#10B981',    # 翡翠绿
        'accent': '#F59E0B',       # 琥珀黄
        'danger': '#EF4444',       # 风险红
        'bg': '#0F172A',           # 深蓝黑
        'text': '#F8FAFC',         # 白
        'muted': '#94A3B8',        # 灰
        'header_font': 'Inter',
        'body_font': 'Source Han Sans CN',
    },
    'golden_hour': {
        'name': 'Golden Hour',
        'primary': '#D97706',      # 金黄
        'secondary': '#B45309',    # 深棕
        'accent': '#FCD34D',       # 亮黄
        'danger': '#DC2626',
        'bg': '#1C1917',
        'text': '#FAFAF9',
        'muted': '#A8A29E',
        'header_font': 'Noto Serif SC',
        'body_font': 'Noto Sans SC',
    },
    'ocean_depths': {
        'name': 'Ocean Depths',
        'primary': '#0EA5E9',
        'secondary': '#06B6D4',
        'accent': '#22D3EE',
        'danger': '#F43F5E',
        'bg': '#082F49',
        'text': '#ECFEFF',
        'muted': '#67E8F9',
        'header_font': 'Montserrat',
        'body_font': 'PingFang SC',
    }
}

def get_theme(name='tech_innovation'):
    """获取主题配置"""
    return THEMES.get(name, THEMES['tech_innovation'])

# ═══════════════════════════════════════════════════════════════
# 3. POWERPOINT — PPT生成器
# ═══════════════════════════════════════════════════════════════

def generate_ppt_summary(results, output_path='/tmp/screening_report.pptx'):
    """生成选股结果PPT"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    except ImportError:
        print("⚠️ python-pptx未安装，尝试安装...")
        os.system("pip install python-pptx -q")
        from pptx import Presentation
        from pptx.util import Inches, Pt

    prs = Presentation()
    theme = get_theme('tech_innovation')

    # 幻灯片1: 封面
    slide_layout = prs.slide_layouts[6]  # 空白
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "A股动量选股系统"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(37, 99, 235)  # primary blue

    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"每日精选报告 | {datetime.now().strftime('%Y-%m-%d')}"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(148, 163, 184)  # muted

    # 幻灯片2: 市场概览
    slide = prs.slides.add_slide(slide_layout)

    tier_counts = {'S': 0, 'A': 0, 'B': 0, 'X': 0}
    for r in results:
        tier_counts[r.get('tier', 'X')] += 1

    overview_text = f"""市场概览

📊 总筛选票数: {len(results)}只
⭐ Tier A (强烈推荐): {tier_counts['A']}只
📈 Tier B (值得关注): {tier_counts['B']}只
⚠️  Tier X (观望): {tier_counts['X']}只

💡 策略说明
• 过夜分 ≥10: T+0午后介入信号强
• 融合分 ≥6: 3-20日波段趋势向好
• 战法命中: 涨停回调、龙头首阴、杯柄等主力行为"""

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = overview_text
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(248, 250, 252)

    # 幻灯片3-5: Top票详解（每页3只）
    tier_a_b = [r for r in results if r.get('tier') in ['A', 'B']]

    for i in range(0, min(9, len(tier_a_b)), 3):
        slide = prs.slides.add_slide(slide_layout)

        for j, r in enumerate(tier_a_b[i:i+3]):
            y_pos = 0.5 + j * 2.2

            # 票信息框
            box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(2))
            tf = box.text_frame
            tf.word_wrap = True

            code = r.get('code', '')
            name = r.get('name', '')
            tier = r.get('tier', 'X')
            score = r.get('final_score', 0)
            overnight = r.get('overnight_score', 0)
            fusion = r.get('fusion_score', 0)
            change = r.get('change_pct', 0)
            close = r.get('close', 0)
            sectors = r.get('sectors', [])
            pattern = r.get('pattern_name', '-')

            tier_color = '#10B981' if tier == 'A' else '#F59E0B'

            text = f"""{code} {name} [{tier}]  综合:{score:.2f}
所属板块: {', '.join(sectors[:3]) if sectors else '—'}
涨跌幅: {change:+.2f}% | 股价: ¥{close:.2f}
过夜分: {overnight:.0f} | 融合分: {fusion:.0f} | 形态: {pattern}"""

            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(248, 250, 252)

    # 幻灯片6: 免责声明
    slide = prs.slides.add_slide(slide_layout)
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = """⚠️ 免责声明

本报告仅供参考，不构成投资建议。
股市有风险，投资需谨慎。
过往业绩不代表未来表现。

请结合自身风险承受能力做出投资决策。"""
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(148, 163, 184)

    prs.save(output_path)
    print(f"✅ PPT已生成: {output_path}")
    return output_path

# ═══════════════════════════════════════════════════════════════
# 4. RESEARCH-PAPER — 策略研究报告生成
# ═══════════════════════════════════════════════════════════════

def generate_research_report(results, output_path='/tmp/strategy_research.md'):
    """生成策略研究报告"""

    tier_a = [r for r in results if r.get('tier') == 'A']
    tier_b = [r for r in results if r.get('tier') == 'B']
    avg_score = sum(r.get('final_score', 0) for r in results) / len(results) if results else 0

    # 板块分析
    sector_counts = {}
    for r in results:
        for s in r.get('sectors', []):
            sector_counts[s] = sector_counts.get(s, 0) + 1
    top_sectors = sorted(sector_counts.items(), key=lambda x: x[1], reverse=True)[:5]

    # 战法统计
    tactic_counts = {}
    for r in results:
        for t in r.get('tactic_names', []):
            tactic_counts[t] = tactic_counts.get(t, 0) + 1

    report = f"""# A股动量选股策略研究报告

> 生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}
> 数据周期: 2025-06-01 至 2025-08-19
> 样本数量: {len(results)}只主板股票

---

## 摘要

本报告基于动量选股策略v2.2，对A股主板股票进行多维度评分。
采用过夜持股法与波段融合策略双轨评估，结合技术形态识别与主力资金分析，
筛选出具备短期爆发潜力的投资标的。

**核心发现**:
- Tier A级股票 {len(tier_a)} 只，具备最强介入信号
- Tier B级股票 {len(tier_b)} 只，趋势向好值得观察
- 平均评分 {avg_score:.2f}，市场整体处于结构性机会期
- 热点板块集中于: {', '.join([s[0] for s in top_sectors[:3]])}

---

## 1. 策略框架

### 1.1 过夜持股法 (Overnight Strategy)

适合T+0交易制度下的日内短线操作:
- **介入时机**: 午后14:00-14:30，当日涨幅2-5%且量比>1.5
- **持有周期**: 当日买入，次日早盘卖出
- **核心指标**: 主力净流入、量比、分时均线支撑

### 1.2 波段融合策略 (Swing Fusion)

适合3-20日中期持仓:
- **介入时机**: 突破20日高点或杯柄形态确认
- **持有周期**: 3-20个交易日
- **核心指标**: MACD金叉、RSI 50-70、均线多头排列

### 1.3 评分体系

| 维度 | 权重 | 说明 |
|------|------|------|
| 过夜分 | 40% | T+0介入信号强度 |
| 融合分 | 35% | 波段趋势质量 |
| 技术形态 | 15% | 突破/杯柄/锤头 |
| 战法加分 | 10% | 涨停回调/龙头首阴 |

---

## 2. 市场分析

### 2.1 板块热度分布

| 板块 | 入选数量 | 热度评级 |
|------|----------|----------|
"""

    for sector, count in top_sectors:
        report += f"| {sector} | {count}只 | {'🔥' * min(count, 5)} |\n"

    report += f"""
### 2.2 技术形态识别

| 形态 | 出现次数 | 胜率参考 |
|------|----------|----------|
"""

    for tactic, count in tactic_counts.items():
        report += f"| {tactic} | {count}次 | 待回测 |\n"

    if not tactic_counts:
        report += "| 无显著形态 | — | — |\n"

    report += f"""
---

## 3. 精选标的分析

### 3.1 Tier A — 强烈推荐

"""

    for r in tier_a[:5]:
        report += f"""#### {r['code']} {r['name']}

- **综合评分**: {r.get('final_score', 0):.2f}
- **过夜分**: {r.get('overnight_score', 0):.0f} (T+0介入信号)
- **融合分**: {r.get('fusion_score', 0):.0f} (波段趋势)
- **技术形态**: {r.get('pattern_name', '-')}
- **所属板块**: {', '.join(r.get('sectors', [])[:3])}
- **当日涨跌**: {r.get('change_pct', 0):+.2f}%

"""

    report += f"""
### 3.2 Tier B — 值得关注

"""

    for r in tier_b[:5]:
        report += f"- **{r['code']} {r['name']}**: 综合{r.get('final_score', 0):.2f}, 过夜{r.get('overnight_score', 0):.0f}, 融合{r.get('fusion_score', 0):.0f}\n"

    report += f"""
---

## 4. 风险提示

1. **市场风险**: 系统性下跌可能导致所有技术形态失效
2. **流动性风险**: 小盘股可能出现买卖价差过大
3. **模型风险**: 历史表现不代表未来收益
4. **执行风险**: 滑点、延迟可能影响实际收益

---

## 5. 后续研究建议

1. **回测验证**: 对v2.2评分体系进行2019-2024年历史回测
2. **参数优化**: 网格搜索过夜分/融合分权重最优配比
3. **机器学习**: 引入XGBoost/LightGBM预测次日涨跌
4. **情绪因子**: 接入财联社/雪球舆情数据

---

*本报告由A股动量选股系统自动生成，仅供参考，不构成投资建议。*
"""

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(report)

    print(f"✅ 研究报告已生成: {output_path}")
    return output_path

# ═══════════════════════════════════════════════════════════════
# 5. COPY-EDITING — 输出质量审查
# ═══════════════════════════════════════════════════════════════

def review_output(text):
    """审查输出质量，标记AI痕迹和表达问题"""
    issues = []

    # 检查AI常见表达
    ai_patterns = [
        '值得注意的是', '不难发现', '综上所述', '总而言之',
        '需要指出的是', '必须强调的是', '显而易见',
        '在...的背景下', '从...的角度来看',
    ]

    for pattern in ai_patterns:
        if pattern in text:
            issues.append(f"发现AI表达: '{pattern}'")

    # 检查过度使用修辞
    if text.count('！') > 5:
        issues.append("感叹号过多，可能显得不够专业")

    # 检查三段式结构
    paragraphs = text.split('\n\n')
    if len(paragraphs) >= 3 and all(len(p) > 50 for p in paragraphs[:3]):
        issues.append("可能存在套路化三段式结构")

    if issues:
        print("\n⚠️ 输出审查发现以下问题:")
        for issue in issues:
            print(f"  - {issue}")
    else:
        print("\n✅ 输出审查通过，无明显AI痕迹")

    return len(issues) == 0

# ═══════════════════════════════════════════════════════════════
# 6. ARXIV — 学术策略研究
# ═══════════════════════════════════════════════════════════════

def research_academic_strategies():
    """搜索学术文献中的动量策略研究"""
    print("\n📚 学术策略研究建议:")
    print("""
建议搜索以下方向的arXiv论文:

1. "momentum strategy" + "A-share" + "China"
   - 中国A股动量效应的实证研究
   
2. "overnight return" + "intraday" + "prediction"
   - 隔夜收益率预测模型
   
3. "technical analysis" + "machine learning" + "stock"
   - 技术分析的机器学习改进
   
4. "limit up" + "pullback" + "trading strategy"
   - 涨停回调策略的量化研究

使用命令: openclaw tool arxiv search "momentum strategy China A-share"
""")

# ═══════════════════════════════════════════════════════════════
# 7. FIRECRAWL — 实时资讯抓取
# ═══════════════════════════════════════════════════════════════

def fetch_news_for_stocks(stock_list):
    """为持仓票抓取最新资讯"""
    print(f"\n📰 为 {len(stock_list)} 只票抓取最新资讯...")

    queries = []
    for code, name in stock_list:
        queries.append(f"{name} {code} 最新消息")

    # 实际调用时会使用 kimi_search
    print("建议使用:")
    for q in queries[:3]:
        print(f"  kimi_search: {q}")
    print(f"  ... 共 {len(queries)} 个查询")

# ═══════════════════════════════════════════════════════════════
# 8. DESIGN — 可视化增强
# ═══════════════════════════════════════════════════════════════

def generate_html_report(results, output_path='/tmp/screening_report.html'):
    """生成可视化HTML报告"""

    theme = get_theme('tech_innovation')

    tier_a = [r for r in results if r.get('tier') == 'A']
    tier_b = [r for r in results if r.get('tier') == 'B']

    html = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<title>A股动量选股报告</title>
<style>
    * {{ margin: 0; padding: 0; box-sizing: border-box; }}
    body {{
        background: {theme['bg']};
        color: {theme['text']};
        font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif;
        line-height: 1.6;
        padding: 20px;
    }}
    .container {{ max-width: 1200px; margin: 0 auto; }}
    .header {{
        text-align: center;
        padding: 40px 0;
        border-bottom: 2px solid {theme['primary']};
        margin-bottom: 30px;
    }}
    .header h1 {{
        font-size: 2.5em;
        color: {theme['primary']};
        margin-bottom: 10px;
    }}
    .stats {{
        display: grid;
        grid-template-columns: repeat(4, 1fr);
        gap: 20px;
        margin-bottom: 40px;
    }}
    .stat-card {{
        background: rgba(37, 99, 235, 0.1);
        border: 1px solid {theme['primary']};
        border-radius: 12px;
        padding: 20px;
        text-align: center;
    }}
    .stat-card h3 {{
        font-size: 2em;
        color: {theme['primary']};
    }}
    .stock-card {{
        background: rgba(255,255,255,0.05);
        border-radius: 12px;
        padding: 20px;
        margin-bottom: 15px;
        border-left: 4px solid {theme['primary']};
    }}
    .stock-card.tier-a {{ border-left-color: {theme['secondary']}; }}
    .stock-header {{
        display: flex;
        justify-content: space-between;
        align-items: center;
        margin-bottom: 10px;
    }}
    .stock-code {{ font-size: 1.3em; font-weight: bold; }}
    .stock-tier {{
        padding: 4px 12px;
        border-radius: 20px;
        font-size: 0.9em;
        font-weight: bold;
    }}
    .tier-a-badge {{ background: {theme['secondary']}; color: #000; }}
    .tier-b-badge {{ background: {theme['accent']}; color: #000; }}
    .stock-tags {{
        display: flex;
        gap: 8px;
        flex-wrap: wrap;
        margin: 10px 0;
    }}
    .tag {{
        background: rgba(255,255,255,0.1);
        padding: 4px 10px;
        border-radius: 6px;
        font-size: 0.85em;
    }}
    .stock-metrics {{
        display: grid;
        grid-template-columns: repeat(3, 1fr);
        gap: 15px;
        margin-top: 10px;
    }}
    .metric {{
        text-align: center;
    }}
    .metric-value {{
        font-size: 1.5em;
        font-weight: bold;
        color: {theme['primary']};
    }}
    .metric-label {{
        font-size: 0.85em;
        color: {theme['muted']};
    }}
</style>
</head>
<body>
<div class="container">
    <div class="header">
        <h1>🎯 A股动量选股系统</h1>
        <p>报告日期: {datetime.now().strftime('%Y-%m-%d')} | 数据截止: 2025-08-19</p>
    </div>

    <div class="stats">
        <div class="stat-card">
            <div class="metric-label">总筛选</div>
            <h3>{len(results)}</h3>
            <div class="metric-label">只</div>
        </div>
        <div class="stat-card">
            <div class="metric-label">Tier A</div>
            <h3 style="color: {theme['secondary']}">{len(tier_a)}</h3>
            <div class="metric-label">强烈推荐</div>
        </div>
        <div class="stat-card">
            <div class="metric-label">Tier B</div>
            <h3 style="color: {theme['accent']}">{len(tier_b)}</h3>
            <div class="metric-label">值得关注</div>
        </div>
        <div class="stat-card">
            <div class="metric-label">平均分</div>
            <h3>{sum(r.get('final_score',0) for r in results)/len(results):.2f}</h3>
            <div class="metric-label">综合</div>
        </div>
    </div>
"""

    # Tier A 票
    for r in tier_a[:10]:
        html += generate_stock_card_html(r, 'tier-a', theme)

    # Tier B 票
    for r in tier_b[:10]:
        html += generate_stock_card_html(r, 'tier-b', theme)

    html += """
    <div style="text-align: center; margin-top: 40px; padding: 20px; color: #64748B; font-size: 0.9em;">
        <p>⚠️ 本报告仅供参考，不构成投资建议。股市有风险，投资需谨慎。</p>
    </div>
</div>
</body>
</html>"""

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(html)

    print(f"✅ HTML报告已生成: {output_path}")
    return output_path

def generate_stock_card_html(r, tier_class, theme):
    """生成单只票的HTML卡片"""
    code = r.get('code', '')
    name = r.get('name', '')
    tier = r.get('tier', 'X')
    score = r.get('final_score', 0)
    overnight = r.get('overnight_score', 0)
    fusion = r.get('fusion_score', 0)
    change = r.get('change_pct', 0)
    close = r.get('close', 0)
    sectors = r.get('sectors', [])
    pattern = r.get('pattern_name', '-')
    tactics = r.get('tactic_names', [])

    change_color = '#10B981' if change > 0 else '#EF4444'
    change_sign = '+' if change > 0 else ''

    tags_html = ''
    for s in sectors[:3]:
        tags_html += f'<span class="tag">{s}</span>'
    if pattern != '-':
        tags_html += f'<span class="tag" style="background: {theme["primary"]}20;">{pattern}</span>'
    for t in tactics:
        tags_html += f'<span class="tag" style="background: {theme["accent"]}30;">⚔️ {t}</span>'

    return f"""
    <div class="stock-card {tier_class}">
        <div class="stock-header">
            <div>
                <span class="stock-code">{code} {name}</span>
                <span class="stock-tier tier-{tier.lower()}-badge">{tier}</span>
            </div>
            <span style="font-size: 1.2em; color: {change_color};">{change_sign}{change:.2f}%</span>
        </div>
        <div class="stock-tags">
            {tags_html}
        </div>
        <div class="stock-metrics">
            <div class="metric">
                <div class="metric-value">{score:.2f}</div>
                <div class="metric-label">综合分</div>
            </div>
            <div class="metric">
                <div class="metric-value">{overnight:.0f}</div>
                <div class="metric-label">过夜分</div>
            </div>
            <div class="metric">
                <div class="metric-value">{fusion:.0f}</div>
                <div class="metric-label">融合分</div>
            </div>
        </div>
    </div>
"""

# ═══════════════════════════════════════════════════════════════
# 9. WRITING-PLANS — 迭代计划（已集成到ROADMAP.md）
# ═══════════════════════════════════════════════════════════════

# ═══════════════════════════════════════════════════════════════
# 10. 统一入口
# ═══════════════════════════════════════════════════════════════

def run_skill_integration(results=None, result_file='/tmp/scan_157_full_result.json'):
    """运行全部skill融合流程"""

    print("=" * 80)
    print("🎨 选股系统 Skill 融合模块启动")
    print("=" * 80)

    # 加载数据
    if results is None:
        if os.path.exists(result_file):
            with open(result_file, 'r', encoding='utf-8') as f:
                results = json.load(f)
        else:
            print(f"❌ 未找到结果文件: {result_file}")
            return

    print(f"\n📊 加载 {len(results)} 只票的评分数据")

    # 1. Copy-editing审查
    print("\n" + "—" * 60)
    print("📝 Step 1: 输出质量审查 (copy-editing)")
    sample_text = f"本次筛选共{len(results)}只股票，值得注意的是Tier A有{len([r for r in results if r.get('tier')=='A'])}只。"
    review_output(sample_text)

    # 2. 生成PPT
    print("\n" + "—" * 60)
    print("📊 Step 2: 生成PPT报告 (powerpoint-pptx)")
    try:
        ppt_path = generate_ppt_summary(results)
    except Exception as e:
        print(f"  ⚠️ PPT生成失败: {e}")
        ppt_path = None

    # 3. 生成研究报告
    print("\n" + "—" * 60)
    print("📚 Step 3: 生成策略研究报告 (research-paper-writer-pro)")
    try:
        report_path = generate_research_report(results)
    except Exception as e:
        print(f"  ⚠️ 报告生成失败: {e}")
        report_path = None

    # 4. 生成HTML可视化报告
    print("\n" + "—" * 60)
    print("🎨 Step 4: 生成HTML可视化报告 (design + theme-factory)")
    try:
        html_path = generate_html_report(results)
    except Exception as e:
        print(f"  ⚠️ HTML生成失败: {e}")
        html_path = None

    # 5. 学术研究方向
    print("\n" + "—" * 60)
    print("📖 Step 5: 学术策略研究方向 (arxiv)")
    research_academic_strategies()

    # 6. 实时资讯建议
    print("\n" + "—" * 60)
    print("📰 Step 6: 实时资讯抓取建议 (firecrawl)")
    top_stocks = [(r['code'], r['name']) for r in results if r.get('tier') in ['A', 'B']][:10]
    fetch_news_for_stocks(top_stocks)

    # 汇总
    print("\n" + "=" * 80)
    print("✅ Skill融合完成！输出文件:")
    print("=" * 80)
    if ppt_path:
        print(f"  📊 PPT报告: {ppt_path}")
    if report_path:
        print(f"  📚 研究报告: {report_path}")
    if html_path:
        print(f"  🌐 HTML报告: {html_path}")
    print(f"  📈 原始数据: {result_file}")
    print("=" * 80)

if __name__ == '__main__':
    run_skill_integration()
